from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # Import RGBColor
from pptx.enum.text import PP_ALIGN  # Import PP_ALIGN

# Define the PDF and PPT file paths
pdf_file_path = './NguyenDinhThaiCV.pdf'
pptx_file_path = 'output.pptx'

# Initialize a new PowerPoint presentation
ppt = Presentation()

# Create a custom slide layout
slide_layout = ppt.slide_layouts[5]  # Use the blank slide layout (5)

# Open the PDF file with PyPDF2
pdf_reader = PdfReader(pdf_file_path)

# Iterate through each page in the PDF
for page_number in range(len(pdf_reader.pages)):
    # Create a new slide in the PowerPoint presentation
    slide = ppt.slides.add_slide(slide_layout)

    # Extract text from the PDF page
    page = pdf_reader.pages[page_number]
    text = page.extract_text()

    # Add the extracted text to the slide
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.alignment = PP_ALIGN.LEFT  # Use PP_ALIGN to set text alignment

    # Set font properties (Times New Roman, white font color, black background)
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Times New Roman'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255, 255, 255)  # White font color using RGBColor
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background color using RGBColor

# Save the PowerPoint presentation
ppt.save(pptx_file_path)

print(f'PDF converted to PPT: {pptx_file_path}')